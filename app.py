import streamlit as st
import fitz  # PyMuPDF
from openai import OpenAI
import os
import docx
from docx import Document
from docx.shared import Pt
import io
from pptx import Presentation
import tempfile
from typing import List, Tuple
import subprocess
import concurrent.futures
from tenacity import retry, wait_random_exponential, stop_after_attempt

# Streamlit app title
st.title("Document Translator with Back Translation and Evaluation")

# User password for page protection
def check_password():
    """Function to check user password to protect the page"""
    password = st.text_input("Enter Password", type="password")
    if password == st.secrets["password"]:
        return True
    elif password:
        st.error("Incorrect password")
    return False

# Text extraction functions
def extract_text_from_pdf(file_stream) -> List[Tuple[int, str]]:
    """Extract text from PDF using PyMuPDF, return list of (page_num, text)"""
    doc = fitz.open(stream=file_stream.read(), filetype="pdf")
    pages = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            pages.append((page_num, text))
    return pages

# Updated extract_text_from_docx function
def extract_text_from_docx(file_stream) -> List[Tuple[int, str]]:
    """Extract text from DOCX using python-docx, return list of (page_num, text)"""
    doc = Document(file_stream)
    full_text = []
    page_texts = []
    page_num = 1

    for para in doc.paragraphs:
        full_text.append(para.text)
        # Heuristic: Assume a new page starts after a certain number of paragraphs
        if len(full_text) >= 10:  # Adjust this number based on typical page length
            page_texts.append((page_num, '\n'.join(full_text).strip()))
            full_text = []
            page_num += 1

    # Add any remaining text as the last page
    if full_text:
        page_texts.append((page_num, '\n'.join(full_text).strip()))

    return page_texts

def extract_text_from_pptx(file_stream) -> List[Tuple[int, str]]:
    """Extract text from PPTX using python-pptx, return list of (slide_num, text)"""
    prs = Presentation(file_stream)
    slides = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        slide_text = slide_text.strip()
        if slide_text:
            slides.append((slide_num, slide_text))
    return slides

def extract_text(uploaded_file) -> Tuple[str, List[Tuple[int, str]]]:
    """Determine file type and extract text accordingly, return file type and list of (unit_num, text)"""
    file_extension = os.path.splitext(uploaded_file.name)[1].lower()
    if file_extension == ".pdf":
        return ("PDF", extract_text_from_pdf(uploaded_file))
    elif file_extension == ".docx":
        return ("DOCX", extract_text_from_docx(uploaded_file))
    elif file_extension == ".pptx":
        return ("PPTX", extract_text_from_pptx(uploaded_file))
    else:
        return (None, [])

def initialize_openai_client(api_key: str) -> OpenAI:
    """Initialize OpenAI client"""
    return OpenAI(api_key=api_key)

@retry(wait=wait_random_exponential(min=1, max=60), stop=stop_after_attempt(6))
def openai_chatbot(client, system_content, user_content, temperature=0.3):
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": system_content},
            {"role": "user", "content": user_content}
        ],
        temperature=temperature,
        max_tokens=4096
    )
    return response.choices[0].message.content.strip()

def parallel_process_openai_chatbot(client, system_content, user_contents):
    input_data = [(client, system_content, user_content) for user_content in user_contents]
    # Progress bar
    progress_bar = st.progress(0)
    progress_text = st.empty()
    with concurrent.futures.ThreadPoolExecutor(max_workers=20) as executor:
        future_to_index = {
            executor.submit(openai_chatbot, client, system_content, user_content): idx
            for idx, (client, system_content, user_content) in enumerate(input_data)
        }
        results = [None] * len(user_contents)
        total = len(user_contents)
        for count, future in enumerate(concurrent.futures.as_completed(future_to_index), start=1):
            index = future_to_index[future]
            try:
                result = future.result()
            except Exception as exc:
                st.write(f'Generated an exception: {exc}')
            else:
                results[index] = result
                progress = count / total
                progress_bar.progress(progress)
                progress_text.text(f"Processing unit {count} of {total}...")
    return results

def sanitize_text(text: str) -> str:
    """Remove non-XML compatible characters from the text."""
    return ''.join(c for c in text if c.isprintable())

def add_markdown_paragraph(doc, text):
    """Add a paragraph to the document with markdown-style bolding and respect line breaks."""
    parts = text.split('\n')
    for line in parts:
        paragraph = doc.add_paragraph()
        bold_parts = line.split('**')
        for i, part in enumerate(bold_parts):
            run = paragraph.add_run(part)
            if i % 2 == 1:  # Odd indices are between pairs of asterisks
                run.bold = True

def main():
    if check_password():
        # Initialize OpenAI client
        api_key = st.secrets["OPENAI_API_KEY"]  # Ensure you have this key in your Streamlit secrets
        client = initialize_openai_client(api_key)

        # File uploader to allow users to upload multiple file types
        uploaded_file = st.file_uploader(
            "Upload a Document",
            type=["pdf", "docx", "pptx"]
        )

        # Input for selecting the language for translation
        language = st.text_input("Select language for translation", "Spanish")

        time_warning = "Please note: processing each unit (page/slide/section) may take several seconds (translation, back translation, and evaluation)."

        st.write(time_warning)

        # Button to start the translation process
        if uploaded_file and language:
            if st.button("Start"):
                with st.spinner("Extracting text from the document..."):
                    file_type, units = extract_text(uploaded_file)
                
                if not units:
                    st.error("Unsupported file type or failed to extract text.")
                    return

                system_text_translate = f"You are an expert at translating languages. Translate the text into {language}."
                user_contents_translate = [f"""
                # Instructions

                Translate the below text to {language}.

                # Text: 

                {unit[1]}

                """ for unit in units]

                # Process translations in parallel
                st.write("# Step 1: Translation")
                translated_texts = parallel_process_openai_chatbot(client, system_text_translate, user_contents_translate)

                # Prepare for back translation
                system_text_back_translate = "You are an expert at translating languages. A user will provide text in a foreign language, and you will translate it back to English."
                user_contents_back_translate = [f"""
                # Instructions

                Translate the below text back to English.

                # Text: 

                {translated_text}
                """ for translated_text in translated_texts]

                # Process back translations in parallel
                st.write("# Step 2: Back Translation (To English)")
                back_translated_texts = parallel_process_openai_chatbot(client, system_text_back_translate, user_contents_back_translate)

                # Prepare user contents for evaluation
                system_text_evaluate = """
                Below I have an original text in English, and a back translated text in English. Provide the following: 
                - A bulleted list that briefly details how accurate the back translation is on tone, voice, style, word choice (1 bullet for each of these)
                - Rewrite the original text and add bolding where the back translation uses a different word
                """
                user_contents_evaluate = [f"""
                # Original Text:

                {unit_text}

                # Back Translated Text:

                {back_translated_text}
                """ for (unit_text, back_translated_text) in zip([unit[1] for unit in units], back_translated_texts)]

                # Process evaluations in parallel
                st.write("# Step 3: Compare Original + Back Translations")
                evaluations = parallel_process_openai_chatbot(client, system_text_evaluate, user_contents_evaluate)

                # Initialize Word Document
                docx_document = Document()
                docx_document.add_heading("Document Translation Report", 0)
                docx_document.add_paragraph(f"Original File: {uploaded_file.name}")
                docx_document.add_paragraph(f"Translated Language: {language}")
                docx_document.add_paragraph(f"Processed File Type: {file_type}")
                docx_document.add_page_break()

                # Loop through each unit and process
                for idx, (unit_num, unit_text) in enumerate(units, start=1):

                    # Add unit heading
                    if file_type in ["PDF", "DOC", "DOCX"]:
                        unit_heading = f"Section {unit_num}"
                    elif file_type in ["PPT", "PPTX"]:
                        unit_heading = f"Slide {unit_num}"
                    else:
                        unit_heading = f"Unit {unit_num}"

                    # Add to Word Document

                    # Add original translation
                    sanitized_unit_text = sanitize_text(unit_text)
                    docx_document.add_heading(f"{unit_heading} of {len(units)}", level=2)
                    docx_document.add_heading("Original Text:", level=3)
                    add_markdown_paragraph(docx_document, sanitized_unit_text)

                    # Display and add Translated Text
                    translated_text = translated_texts[idx - 1]  # Use precomputed translated text
                    sanitized_translated_text = sanitize_text(translated_text)
                    docx_document.add_heading("Translated Text:", level=3)
                    add_markdown_paragraph(docx_document, sanitized_translated_text)

                    # Display and add Back Translated Text
                    back_translated_text = back_translated_texts[idx - 1]  # Use precomputed back-translated text
                    sanitized_back_translated_text = sanitize_text(back_translated_text)
                    docx_document.add_heading("Back Translated Text:", level=3)
                    add_markdown_paragraph(docx_document, sanitized_back_translated_text)

                    # Display and add Evaluation
                    evaluation = evaluations[idx - 1]  # Use precomputed evaluation
                    sanitized_evaluation = sanitize_text(evaluation)
                    docx_document.add_heading("Evaluation of Differences:", level=3)
                    add_markdown_paragraph(docx_document, sanitized_evaluation)

                    docx_document.add_page_break()  # Add a page break after each unit's content

                # After processing all units, provide the download button
                # Save the Word document to a BytesIO buffer
                buffer = io.BytesIO()
                docx_document.save(buffer)
                buffer.seek(0)

                # Provide download button
                st.success("Translation and evaluation completed!")
                st.download_button(
                    label="Download Translation Report",
                    data=buffer,
                    file_name="translation_report.docx",
                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                )

if __name__ == "__main__":
    main()
